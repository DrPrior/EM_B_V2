"""Data ingestion module for populating the Neo4j knowledge graph with embeddings.

This module traverses the project_data directory, reads file contents, chunks text
semantically, generates embeddings via Ollama, and stores everything in Neo4j.
It uses parameterized Cypher queries to prevent injection vulnerabilities.
"""

import re
import sys
import uuid
from pathlib import Path  # type: ignore[import-untyped]

# Third-party imports
from neo4j import Session  # type: ignore[import-untyped]

# Local imports
from src.core.config import settings
from src.database.connection import Neo4jConnection
from src.services.embeddings import embed_and_store_chunk

# Add project root to Python path so local package imports work
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))


# ==========================================
# HELPER FUNCTIONS: FILE I/O
# ==========================================


def discover_files(root_path: str) -> list[str]:
    """Discover all readable files in the directory tree.

    Supports plaintext (.txt, .md), PDF, Word (.docx), and PowerPoint (.pptx) files.

    Args:
        root_path: Root directory to traverse.

    Returns:
        List of absolute file paths.
    """
    supported_extensions = {".txt", ".md", ".pdf", ".docx", ".pptx"}
    files = []

    root = Path(root_path)
    if not root.exists():
        print(f"Warning: Root path does not exist: {root_path}")
        return files

    for file_path in root.rglob("*"):
        if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
            files.append(str(file_path))

    return sorted(files)


def read_file_content(file_path: str) -> str | None:
    """Read text content from a file.

    Supports plaintext, markdown, and PDF files.

    Args:
        file_path: Absolute path to the file.

    Returns:
        File content as string, or None if reading fails.
    """
    try:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix in {".txt", ".md"}:
            with open(path, encoding="utf-8") as f:
                return f.read()

        elif suffix == ".pdf":
            try:
                from pypdf import PdfReader  # type: ignore[import-untyped]

                reader = PdfReader(path)
                text_parts = []
                for page in reader.pages:
                    text_parts.append(page.extract_text())
                return "\n".join(text_parts)
            except ImportError:
                print(f"Warning: pypdf not installed. Skipping PDF: {file_path}")
                return None

        elif suffix == ".docx":
            try:
                import docx  # type: ignore[import-untyped]

                doc = docx.Document(str(path))
                text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n".join(text_parts)
            except ImportError:
                print(f"Warning: python-docx not installed. Skipping DOCX: {file_path}")
                return None

        elif suffix == ".pptx":
            try:
                from pptx import Presentation  # type: ignore[import-untyped]

                prs = Presentation(str(path))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "").strip()
                        if text:
                            text_parts.append(text)
                return "\n".join(text_parts)
            except ImportError:
                print(f"Warning: python-pptx not installed. Skipping PPTX: {file_path}")
                return None

        else:
            print(f"Warning: Unsupported file type: {file_path}")
            return None

    except Exception as e:
        print(f"Error reading file {file_path}: {e}")
        return None


# ==========================================
# HELPER FUNCTIONS: TEXT CHUNKING
# ==========================================


def estimate_tokens(text: str) -> int:
    """Estimate token count using character approximation.

    Approximate: 1 token ≈ 4 characters (empirical for English).

    Args:
        text: Text to estimate.

    Returns:
        Approximate token count.
    """
    return len(text) // 4


def chunk_text(text: str, max_chunk_size: int = 512, overlap: int = 64) -> list[str]:
    """Split text into overlapping chunks with semantic boundaries.

    Splits on sentence endings, newlines before headings/bullets, and blank
    lines so that PowerPoint slides and bullet lists stay coherent. Aims for
    max_chunk_size tokens (approximate), with overlap tokens carried to the
    next chunk.

    Args:
        text: Text to chunk.
        max_chunk_size: Target max tokens per chunk (approx).
        overlap: Tokens of overlap between consecutive chunks.

    Returns:
        List of text chunks.
    """
    max_chars = max_chunk_size * 4
    overlap_chars = overlap * 4

    # Split on:
    #   - sentence endings (.!?) followed by whitespace
    #   - blank lines (paragraph breaks)
    #   - newlines before a capital letter or common bullet markers (-, •, *, digits)
    segments = re.split(
        r"(?<=[.!?])\s+|(?:\r?\n){2,}|\r?\n(?=[A-Z\-•\*\d])",
        text,
    )

    chunks: list[str] = []
    current_chunk = ""

    for segment in segments:
        segment = segment.strip()
        if not segment:
            continue

        test_chunk = current_chunk + " " + segment if current_chunk else segment

        if len(test_chunk) <= max_chars:
            current_chunk = test_chunk
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())

            if len(current_chunk) > overlap_chars:
                overlap_text = current_chunk[-overlap_chars:]
                current_chunk = overlap_text + " " + segment
            else:
                current_chunk = segment

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks


# ==========================================
# HELPER FUNCTIONS: NEO4J OPERATIONS
# ==========================================


def is_file_embedded(session: Session, file_path: str) -> bool:
    """Return True if the file already has at least one embedded Chunk in Neo4j.

    Args:
        session: Neo4j session.
        file_path: Absolute path of the file to check.

    Returns:
        True if the Document exists and has at least one Chunk with a stored embedding.
    """
    result = session.execute_read(
        lambda tx: tx.run(
            """
            MATCH (doc:Document {filepath: $filepath})-[:HAS_CHUNK]->(c:Chunk)
            WHERE c.embedding IS NOT NULL
            RETURN count(c) > 0 AS embedded
            """,
            filepath=file_path,
        ).single()
    )
    return bool(result["embedded"]) if result else False


def create_directory_node(session: Session, path: str, name: str) -> None:
    """Create or merge a Directory node.

    Args:
        session: Neo4j session.
        path: Full directory path.
        name: Directory name.
    """
    query = """
    MERGE (d:Directory {path: $path})
    ON CREATE SET d.name = $name
    """
    session.execute_write(lambda tx: tx.run(query, path=path, name=name))


def _link_directories(session: Session, parent_path: str, child_path: str) -> None:
    """Create a CONTAINS_DIR relationship between two Directory nodes.

    Args:
        session: Neo4j session.
        parent_path: Full path of the parent directory.
        child_path: Full path of the child directory.
    """
    query = """
    MATCH (parent:Directory {path: $parent_path})
    MATCH (child:Directory {path: $child_path})
    MERGE (parent)-[:CONTAINS_DIR]->(child)
    """
    session.execute_write(
        lambda tx: tx.run(query, parent_path=parent_path, child_path=child_path)
    )


def create_directory_hierarchy(session: Session, file_path: str, data_root: str) -> str:
    """Create Directory nodes for every ancestor from data_root to the file's
    parent, linking each level with a CONTAINS_DIR relationship.

    Args:
        session: Neo4j session.
        file_path: Absolute path to the file being ingested.
        data_root: Root ingestion directory (inclusive upper bound).

    Returns:
        The path string of the file's immediate parent directory.
    """
    root = Path(data_root)
    dirs: list[Path] = []
    current = Path(file_path).parent
    while True:
        dirs.append(current)
        if current == root:
            break
        current = current.parent

    dirs.reverse()  # root first, deepest last

    for i, dir_path in enumerate(dirs):
        create_directory_node(session, str(dir_path), dir_path.name)
        if i > 0:
            _link_directories(session, str(dirs[i - 1]), str(dir_path))

    return str(Path(file_path).parent)


def create_document_node(
    session: Session, dir_path: str, filepath: str, filename: str, extension: str
) -> None:
    """Create or merge a Document node and link to parent Directory.

    Args:
        session: Neo4j session.
        dir_path: Parent directory path.
        filepath: Full file path.
        filename: File name (with extension).
        extension: File extension (e.g., ".txt").
    """
    query = """
    MATCH (d:Directory {path: $dir_path})
    MERGE (doc:Document {filepath: $filepath})
    ON CREATE SET doc.filename = $filename, doc.extension = $extension
    MERGE (d)-[:CONTAINS_FILE]->(doc)
    """
    session.execute_write(
        lambda tx: tx.run(
            query,
            dir_path=dir_path,
            filepath=filepath,
            filename=filename,
            extension=extension,
        )
    )


def create_chunk_node(
    session: Session, filepath: str, chunk_id: str, text: str, sequence: int
) -> None:
    """Create a Chunk node and link to parent Document.

    Args:
        session: Neo4j session.
        filepath: Parent document filepath.
        chunk_id: Unique chunk identifier.
        text: Chunk text content.
        sequence: Sequence number (0-indexed).
    """
    query = """
    MATCH (doc:Document {filepath: $filepath})
    MERGE (c:Chunk {chunk_id: $chunk_id})
    ON CREATE SET c.text = $text, c.sequence = $sequence
    MERGE (doc)-[:HAS_CHUNK]->(c)
    """
    session.execute_write(
        lambda tx: tx.run(
            query,
            filepath=filepath,
            chunk_id=chunk_id,
            text=text,
            sequence=sequence,
        )
    )


# ==========================================
# PROGRESS DISPLAY
# ==========================================


def _print_progress(
    index: int, total: int, file_path: str, chunks: int, embeddings: int
) -> None:
    """Print a single-file progress line to the terminal.

    Args:
        index: 1-based position of the current file.
        total: Total number of files to process.
        file_path: Absolute path of the file being processed.
        chunks: Number of chunks created for this file.
        embeddings: Number of embeddings stored for this file.
    """
    pct = int(index / total * 100)
    bar_filled = pct // 5
    bar = "#" * bar_filled + "-" * (20 - bar_filled)
    filename = Path(file_path).name
    print(f"  [{bar}] {pct:3d}%  ({index:>{len(str(total))}}/{total})  {filename}")
    print(f"         chunks: {chunks}  |  embeddings: {embeddings}")


def _print_done(stats: dict[str, int]) -> None:
    """Print the ingestion completion summary to the terminal.

    Args:
        stats: Ingestion statistics dictionary returned by ingest_project_data.
    """
    width = 44
    border = "=" * width
    print(f"\n{border}")
    print("  INGESTION COMPLETE")
    print(border)
    print(f"  Files processed  : {stats['files_processed']}")
    print(f"  Chunks created   : {stats['chunks_created']}")
    print(f"  Embeddings stored: {stats['embeddings_stored']}")
    print(f"  Errors           : {stats['errors']}")
    print(f"{border}\n")


# ==========================================
# MAIN INGESTION ORCHESTRATOR
# ==========================================


def ingest_project_data(
    session: Session, data_root: str = "/app/project_data"
) -> dict[str, int]:
    """Ingest all data from project_data into Neo4j with embeddings.

    Args:
        session: Neo4j session (from connection manager).
        data_root: Root path to ingest (default: Docker container path).

    Returns:
        Dictionary with ingestion statistics (files, chunks, embeddings).
    """
    stats = {
        "files_processed": 0,
        "chunks_created": 0,
        "embeddings_stored": 0,
        "errors": 0,
    }

    print(f"Starting data ingestion from: {data_root}")

    # Discover all files
    files = discover_files(data_root)
    total = len(files)
    print(f"Found {total} files to process\n")

    for index, file_path in enumerate(files, start=1):
        path_obj = Path(file_path)
        filename = path_obj.name
        extension = path_obj.suffix

        if is_file_embedded(session, file_path):
            print(f"  ⏭  [{index}/{total}] Already embedded, skipping: {filename}")
            continue

        # Read file content
        content = read_file_content(file_path)
        if content is None:
            print(f"  ⚠️  [{index}/{total}] Skipping (unreadable): {filename}")
            stats["errors"] += 1
            continue

        # Create directory nodes for the full ancestor chain and link them
        try:
            dir_path_str = create_directory_hierarchy(session, file_path, data_root)
        except Exception as e:
            print(f"  ⚠️  [{index}/{total}] Directory hierarchy failed: {e}")
            stats["errors"] += 1
            continue

        # Create document node
        try:
            create_document_node(session, dir_path_str, file_path, filename, extension)
        except Exception as e:
            print(f"  ⚠️  [{index}/{total}] Document node failed: {e}")
            stats["errors"] += 1
            continue

        stats["files_processed"] += 1

        # Chunk the content and generate embeddings
        chunks = chunk_text(
            content, settings.chunk_max_tokens, settings.chunk_overlap_tokens
        )
        file_chunks = 0
        file_embeddings = 0

        for sequence, chunk_text_content in enumerate(chunks):
            chunk_id = str(uuid.uuid4())
            try:
                create_chunk_node(
                    session, file_path, chunk_id, chunk_text_content, sequence
                )
                stats["chunks_created"] += 1
                file_chunks += 1

                embed_and_store_chunk(chunk_text_content, chunk_id, session)
                stats["embeddings_stored"] += 1
                file_embeddings += 1

            except Exception as e:
                print(f"  ⚠️  Chunk {sequence} failed: {e}")
                stats["errors"] += 1

        _print_progress(index, total, file_path, file_chunks, file_embeddings)

    _print_done(stats)
    return stats


if __name__ == "__main__":
    # Standalone execution: initialize connection and ingest
    nc = Neo4jConnection.get_instance()
    try:
        with nc.session() as session:
            ingest_project_data(session)
    finally:
        nc.close()
